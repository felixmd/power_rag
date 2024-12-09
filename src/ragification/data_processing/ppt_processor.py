from typing import Dict, Any, List
import os
from pptx import Presentation
from pathlib import Path

def process_pptx_file(file_path: str) -> List[Dict[str, Any]]:
    """Process a PPTX file and return a list of slide contents with metadata"""
    prs = Presentation(file_path)
    documents = []
    file_name = Path(file_path).stem
    
    for idx, slide in enumerate(prs.slides, 1):
        text_content = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_content.append(shape.text)
        
        if text_content:
            documents.append({
                'original_id': f"{file_name}_slide_{idx}",
                'content': ' '.join(text_content),
                'metadata': {
                    'file_name': file_name,
                    'slide_number': idx,
                    'source': 'pptx'
                }
            })
    
    return documents

def load_pptx_directory(directory_path: str) -> List[Dict[str, Any]]:
    """Load all PPTX files from a directory"""
    all_documents = []
    
    if not os.path.exists(directory_path):
        os.makedirs(directory_path)
        return all_documents
    
    for file in os.listdir(directory_path):
        if file.endswith('.pptx'):
            file_path = os.path.join(directory_path, file)
            try:
                documents = process_pptx_file(file_path)
                all_documents.extend(documents)
            except Exception as e:
                print(f"Error processing {file}: {str(e)}")
    
    return all_documents 